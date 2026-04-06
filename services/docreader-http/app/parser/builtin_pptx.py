"""
builtin_pptx：python-pptx 回退，按幻灯片聚合文本。
"""

from __future__ import annotations

import io
import logging
from typing import Any

from pptx import Presentation

from app.models.document import Document
from app.parser.base_parser import BaseParser

logger = logging.getLogger(__name__)


class PptxBuiltinParser(BaseParser):
    """PptxBuiltinParser 遍历幻灯片形状上的文本。"""

    def parse_into_text(self, content: bytes) -> Document:
        meta: dict[str, Any] = {"parser": "python-pptx"}
        try:
            prs = Presentation(io.BytesIO(content))
        except Exception as e:
            logger.exception("pptx open failed")
            return Document(metadata={**meta, "error": str(e)})

        slides_text: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    t = shape.text.strip()
                    if t:
                        texts.append(t)
            if texts:
                slides_text.append(f"## Slide {idx}\n\n" + "\n\n".join(texts))
        meta["slide_count"] = len(prs.slides)
        md = "\n\n".join(slides_text)
        if not md.strip():
            return Document(metadata=meta)
        return Document(content=md, metadata=meta)
